from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color Palette ──────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x11, 0x17)   # near-black
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)   # vivid purple
PINK       = RGBColor(0xEC, 0x4E, 0x9B)   # hot pink
ORANGE     = RGBColor(0xFF, 0x7B, 0x54)   # warm orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CODE_BG    = RGBColor(0x1E, 0x22, 0x2A)   # dark code panel
CODE_TEXT  = RGBColor(0xA8, 0xFF, 0x78)   # green code text
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)   # sky blue accent

APP_SCREENSHOT = r"C:\Users\Shama\.gemini\antigravity\brain\8f1f30d3-ed7e-4bba-99cd-159004aa20a5\photo_converter_app_main_1776964515403.png"
OUT_PATH = r"C:\Users\Shama\OneDrive\Desktop\PhotoConverter_Presentation.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ─── helpers ───────────────────────────────────────────────
def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def add_code_box(slide, code, l, t, w, h, font_size=9):
    rect = add_rect(slide, l, t, w, h, CODE_BG)
    txb = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.12),
                                    Inches(w-0.3),  Inches(h-0.24))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.size = Pt(font_size)
    run.font.color.rgb = CODE_TEXT
    run.font.name = "Courier New"

def slide_header(slide, title, subtitle=None):
    # purple top bar
    add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
    add_text(slide, title, 0.4, 0.1, 12, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12, 0.5, size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
# big gradient rect (simulate with two rects)
add_rect(s, 0, 0, 13.33, 7.5, RGBColor(0x14, 0x0A, 0x2E))
add_rect(s, 0, 3.5, 13.33, 4.0, RGBColor(0x0D, 0x11, 0x17))

# purple accent line
add_rect(s, 1.5, 2.6, 10.3, 0.06, PURPLE)

add_text(s, "📸  Photo Converter", 1.5, 1.2, 10.3, 1.2,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Art Filter Engine & Background Removal",
         1.5, 2.7, 10.3, 0.7, size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# feature pills
features = [("✏️  Sketch", 1.6), ("✨  Anime", 3.8), ("🎨  Painting", 6.0), ("✂️  Remove BG", 8.2)]
for label, x in features:
    add_rect(s, x, 3.7, 2.0, 0.55, PURPLE)
    add_text(s, label, x+0.1, 3.72, 1.8, 0.5, size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "🌐  https://photo-converter-eight.vercel.app",
         1.5, 5.2, 10.3, 0.6, size=16, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "Built with  React  •  FastAPI  •  OpenCV  •  rembg",
         1.5, 5.9, 10.3, 0.6, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Deployed on  Vercel  +  Render  •  GitHub: Itian-shama/photo-converter",
         1.5, 6.5, 10.3, 0.6, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – App Screenshot
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "App Interface", "Live at: https://photo-converter-eight.vercel.app")
if os.path.exists(APP_SCREENSHOT):
    s.shapes.add_picture(APP_SCREENSHOT, Inches(1.0), Inches(1.3),
                          Inches(11.3), Inches(5.8))
else:
    add_text(s, "[ App Screenshot ]", 1, 1.5, 11, 5, size=24,
             color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – What Does It Do?
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "What Does It Do?", "4 powerful AI-powered photo transformations")

cards = [
    ("✏️", "Pencil Sketch",       "Converts your photo into a realistic hand-drawn pencil sketch using grayscale + gaussian blur + color-dodge blending.",       PURPLE),
    ("✨", "Anime Art",           "Transforms photos into anime/cartoon-style art using edge-preserving filters + adaptive thresholding + saturation boost.",   RGBColor(0x06,0x96,0xDD)),
    ("🎨", "Oil Painting",        "Gives your photo a painterly look using edge-preserving smoothing + saturation boost + sharpening kernel.",                   RGBColor(0x10,0xB9,0x81)),
    ("✂️", "Background Removal", "Automatically detects and removes the background using the rembg AI model (U2Net) leaving a transparent PNG.",              ORANGE),
]

for i, (icon, title, desc, color) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.55
    y = 1.35 + row * 2.8
    add_rect(s, x, y, 6.2, 2.5, RGBColor(0x1A, 0x1F, 0x2E))
    add_rect(s, x, y, 0.12, 2.5, color)
    add_text(s, f"{icon}  {title}", x+0.3, y+0.2, 5.8, 0.55, size=18, bold=True, color=color)
    add_text(s, desc, x+0.3, y+0.75, 5.7, 1.6, size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – Tech Stack
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Tech Stack", "Technologies used to build the app")

add_text(s, "🖥️  FRONTEND", 0.5, 1.1, 5.8, 0.5, size=15, bold=True, color=ACCENT)
fe = [
    ("React 19",     "UI framework with hooks & state management"),
    ("Vite 8",       "Ultra-fast build tool & dev server"),
    ("Vanta.js",     "Animated HALO background effect"),
    ("Vanilla CSS",  "Custom styling with glassmorphism effects"),
    ("Vercel",       "Free global CDN hosting"),
]
for j, (tech, desc) in enumerate(fe):
    y = 1.6 + j * 0.85
    add_rect(s, 0.5, y, 5.8, 0.75, RGBColor(0x1A, 0x1F, 0x2E))
    add_text(s, tech, 0.7, y+0.08, 1.8, 0.55, size=13, bold=True, color=PURPLE)
    add_text(s, desc, 2.5, y+0.08, 3.7, 0.55, size=11, color=LIGHT_GRAY)

add_text(s, "⚙️  BACKEND", 7.0, 1.1, 5.8, 0.5, size=15, bold=True, color=ORANGE)
be = [
    ("FastAPI",               "High-performance Python web framework"),
    ("OpenCV (contrib)",      "Computer vision – sketch, anime, painting"),
    ("rembg + U2Net",         "AI background removal model"),
    ("ONNX Runtime",          "Fast ML inference engine"),
    ("Render.com",            "Free Python server hosting"),
]
for j, (tech, desc) in enumerate(be):
    y = 1.6 + j * 0.85
    add_rect(s, 7.0, y, 5.8, 0.75, RGBColor(0x1A, 0x1F, 0x2E))
    add_text(s, tech, 7.2, y+0.08, 2.2, 0.55, size=13, bold=True, color=ORANGE)
    add_text(s, desc, 9.4, y+0.08, 3.2, 0.55, size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – Project Structure
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Project Structure", "How the code is organized")

structure = """\
📦 photo-converter/
├── 🖥️  start.bat              ← Run this to launch the app locally
│
├── 📁 backend/
│   ├── 🐍 main.py             ← FastAPI server with all 4 endpoints
│   ├── 📄 requirements.txt    ← Python dependencies
│   └── 📄 render.yaml         ← Render deployment config
│
└── 📁 frontend/
    ├── 📄 index.html           ← Entry HTML (loads Vanta, Three.js CDN)
    ├── 📄 package.json         ← Node dependencies
    ├── 📄 vite.config.js       ← Vite dev server config
    ├── 📄 vercel.json          ← Vercel deployment config
    └── 📁 src/
        ├── ⚛️  App.jsx          ← Main React component (UI + API calls)
        ├── 🎨 index.css        ← All styling
        └── 📄 main.jsx         ← React app entry point"""

add_code_box(s, structure, 0.5, 1.1, 12.3, 6.0, font_size=12)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – Backend Code
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Backend Code  —  main.py", "FastAPI endpoints for all 4 image transformations")

code_left = """\
# main.py  – FastAPI Backend
import cv2, numpy as np
from fastapi import FastAPI, File, UploadFile
from fastapi.responses import Response
from fastapi.middleware.cors import CORSMiddleware
from rembg import remove

app = FastAPI()
app.add_middleware(CORSMiddleware,
    allow_origins=["*"], allow_methods=["*"],
    allow_headers=["*"], allow_credentials=True)

@app.post("/api/sketch")
async def convert_to_sketch(file: UploadFile = File(...)):
    img  = cv2.imdecode(np.frombuffer(
               await file.read(), np.uint8), cv2.IMREAD_COLOR)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    inv  = cv2.bitwise_not(gray)
    blur = cv2.GaussianBlur(inv, (21,21), 0)
    sketch = cv2.divide(gray, cv2.bitwise_not(blur), scale=256.0)
    _, out = cv2.imencode('.jpg', sketch)
    return Response(content=out.tobytes(),
                    media_type="image/jpeg")"""

code_right = """\
@app.post("/api/anime")
async def apply_anime(file: UploadFile = File(...)):
    img   = cv2.imdecode(np.frombuffer(
                await file.read(), np.uint8), cv2.IMREAD_COLOR)
    color = cv2.edgePreservingFilter(img, flags=1,
                sigma_s=50, sigma_r=0.4)
    # boost saturation
    hsv = cv2.cvtColor(color,
              cv2.COLOR_BGR2HSV).astype("float32")
    h,s,v = cv2.split(hsv)
    s = np.clip(s*1.25, 0, 255)
    color = cv2.cvtColor(
        cv2.merge([h,s,v]).astype("uint8"),
        cv2.COLOR_HSV2BGR)
    # get outline edges
    gray  = cv2.medianBlur(
        cv2.cvtColor(img, cv2.COLOR_BGR2GRAY), 5)
    edges = cv2.adaptiveThreshold(gray, 255,
        cv2.ADAPTIVE_THRESH_MEAN_C,
        cv2.THRESH_BINARY, 9, 9)
    anime = cv2.bitwise_and(color, color, mask=edges)
    _, out = cv2.imencode('.jpg', anime)
    return Response(out.tobytes(),
                    media_type="image/jpeg")

@app.post("/api/remove-bg")
async def remove_background(file: UploadFile = File(...)):
    result = remove(await file.read())
    return Response(result, media_type="image/png")"""

add_code_box(s, code_left,  0.3, 1.1, 6.3, 6.1, font_size=8)
add_code_box(s, code_right, 6.75, 1.1, 6.3, 6.1, font_size=8)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 – Frontend Code
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Frontend Code  —  App.jsx", "React component that handles UI, uploads & API calls")

code_fe = """\
// App.jsx  – React Frontend
import { useState, useRef, useEffect } from 'react';

function App() {
  const [originalImage, setOriginalImage] = useState(null);
  const [resultImage,   setResultImage]   = useState(null);
  const [isLoading,     setIsLoading]     = useState(false);
  const [activeMode,    setActiveMode]    = useState('sketch');
  const fileInputRef = useRef(null);

  // ── Process image by calling the backend ──
  const processImage = async () => {
    if (!originalImage) return;
    setIsLoading(true);
    setResultImage(null);

    const formData = new FormData();
    formData.append('file', fileInputRef.current.files[0]);

    // Choose endpoint based on selected mode
    const endpoints = {
      sketch:    '/api/sketch',
      anime:     '/api/anime',
      painting:  '/api/painting',
      'remove-bg': '/api/remove-bg',
    };

    const BACKEND_URL =
      import.meta.env.VITE_BACKEND_URL || 'http://localhost:5000';

    try {
      const res = await fetch(
        `${BACKEND_URL}${endpoints[activeMode]}`,
        { method: 'POST', body: formData }
      );
      if (res.ok) {
        const blob = await res.blob();
        setResultImage(URL.createObjectURL(blob));  // show result
      }
    } catch (err) {
      console.error("Error processing image:", err);
    } finally {
      setIsLoading(false);
    }
  };

  // ── JSX: mode buttons + upload area + result display ──
  return (
    <div className="app-container">
      <h1>Creative Studio</h1>
      {/* Mode selector buttons */}
      {['sketch','anime','painting','remove-bg'].map(mode => (
        <div key={mode}
             className={activeMode===mode ? 'active' : ''}
             onClick={() => setActiveMode(mode)}>
          {mode}
        </div>
      ))}
      {/* Upload & Convert */}
      <input type="file" ref={fileInputRef}
             onChange={e => setOriginalImage(
               URL.createObjectURL(e.target.files[0]))} />
      <button onClick={processImage} disabled={isLoading}>
        {isLoading ? 'Processing...' : 'Convert'}
      </button>
      {/* Show original & result side by side */}
      {originalImage && <img src={originalImage} alt="Original" />}
      {resultImage   && <img src={resultImage}   alt="Result"   />}
    </div>
  );
}
export default App;"""

add_code_box(s, code_fe, 0.3, 1.1, 12.7, 6.1, font_size=8)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 – How It Works
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "How It Works", "Step-by-step flow from upload to result")

steps = [
    ("1", "User Uploads Photo",     "Click the upload area → select any JPG or PNG photo from your device.",                              PURPLE),
    ("2", "Choose a Filter",        "Click one of the 4 buttons: Pencil Sketch, Anime Art, Painting, or Remove Background.",              RGBColor(0x06,0x96,0xDD)),
    ("3", "Frontend Sends Request", "React reads the file → creates a FormData object → sends a POST request to the FastAPI backend.",     ORANGE),
    ("4", "Backend Processes",      "OpenCV / rembg applies the chosen transformation on the server and returns the processed image bytes.",RGBColor(0x10,0xB9,0x81)),
    ("5", "Result Displayed",       "React receives the image blob → creates an object URL → displays it beside the original photo.",      PINK),
    ("6", "Download",               "A 'Download Result' link appears so the user can save the transformed image.",                        ACCENT),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.55
    y = 1.35 + row * 1.95
    add_rect(s, x, y, 6.2, 1.75, RGBColor(0x1A,0x1F,0x2E))
    # number circle
    add_rect(s, x+0.12, y+0.35, 0.55, 0.55, color)
    add_text(s, num, x+0.12, y+0.32, 0.55, 0.6, size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+0.8, y+0.18, 5.2, 0.5, size=14, bold=True, color=color)
    add_text(s, desc,  x+0.8, y+0.65, 5.2, 1.0, size=10, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 – Deployment Steps
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Deployment Steps", "How we made the app available to the world for free")

deploy_steps = [
    ("🐙  Step 1 – GitHub",
     "Push Code to GitHub",
     ["git init  (initialize repo)",
      "git add -A  (stage all files)",
      'git commit -m "Initial commit"',
      "git remote add origin https://github.com/Itian-shama/photo-converter.git",
      "git push -u origin main"],
     PURPLE),
    ("⚙️  Step 2 – Render (Backend)",
     "Deploy Python Backend for Free",
     ["Sign up at render.com using GitHub",
      'New → Web Service → select "photo-converter"',
      "Root Directory: backend",
      "Build: pip install -r requirements.txt",
      "Start: uvicorn main:app --host 0.0.0.0 --port $PORT",
      "✅  Live at: photo-converter-backend.onrender.com"],
     ORANGE),
    ("🌐  Step 3 – Vercel (Frontend)",
     "Deploy React Frontend for Free",
     ["Sign up at vercel.com using GitHub",
      'New Project → import "photo-converter" repo',
      "Root Directory: frontend",
      "Env var: VITE_BACKEND_URL = (Render URL)",
      "Click Deploy",
      "✅  Live at: photo-converter-eight.vercel.app"],
     ACCENT),
]

for i, (label, title, bullets, color) in enumerate(deploy_steps):
    x = 0.3 + i * 4.35
    add_rect(s, x, 1.1, 4.1, 6.0, RGBColor(0x1A,0x1F,0x2E))
    add_rect(s, x, 1.1, 4.1, 0.08, color)
    add_text(s, label, x+0.15, 1.15, 3.8, 0.5, size=14, bold=True, color=color)
    add_text(s, title, x+0.15, 1.65, 3.8, 0.55, size=12, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        add_text(s, f"›  {b}", x+0.15, 2.3+j*0.77, 3.8, 0.7, size=9.5, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 – Live Links & Summary
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, 13.33, 7.5, RGBColor(0x14, 0x0A, 0x2E))
slide_header(s, "Live App & Links", "Share your app with anyone in the world!")

links = [
    ("🌐  Frontend (Vercel)",  "https://photo-converter-eight.vercel.app",    PURPLE),
    ("⚙️  Backend  (Render)",  "https://photo-converter-backend.onrender.com", ORANGE),
    ("🐙  GitHub Repo",        "https://github.com/Itian-shama/photo-converter", ACCENT),
]

for i, (label, url, color) in enumerate(links):
    y = 1.5 + i * 1.55
    add_rect(s, 1.5, y, 10.3, 1.2, RGBColor(0x1A,0x1F,0x2E))
    add_rect(s, 1.5, y, 0.12, 1.2, color)
    add_text(s, label, 1.8, y+0.1,  9.8, 0.5, size=17, bold=True, color=color)
    add_text(s, url,   1.8, y+0.6,  9.8, 0.5, size=14, color=WHITE)

add_text(s, "⚠️  Note: The free Render backend sleeps after 15 min inactivity.",
         1.5, 6.2, 10.3, 0.5, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "First request after sleep may take ~30-50 sec to wake up. Subsequent requests are fast!",
         1.5, 6.7, 10.3, 0.5, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 – Thank You
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, 13.33, 7.5, RGBColor(0x14, 0x0A, 0x2E))
add_rect(s, 0, 3.2, 13.33, 0.08, PURPLE)

add_text(s, "🎉", 5.5, 0.8, 2.3, 1.5, size=72, align=PP_ALIGN.CENTER)
add_text(s, "Thank You!", 2.0, 2.2, 9.3, 1.2,
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Photo Converter — Built with ❤️ using React + FastAPI + OpenCV",
         1.5, 3.5, 10.3, 0.7, size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "4 filters  •  AI background removal  •  Deployed globally for free",
         1.5, 4.2, 10.3, 0.7, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "🌐  photo-converter-eight.vercel.app",
         1.5, 5.3, 10.3, 0.7, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "🐙  github.com/Itian-shama/photo-converter",
         1.5, 6.0, 10.3, 0.6, size=14, color=PURPLE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"\nDone! Presentation saved to:\n    {OUT_PATH}\n")
print(f"Total slides: {len(prs.slides)}")
